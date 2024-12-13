from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize a PowerPoint presentation
presentation = Presentation()

# Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Movie Recommendation System"
subtitle.text = "A Comprehensive Analysis of Our Movie Recommendation Approach"

# Slide 2: Genre Metrics Table
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Genre Metrics: Average Ratings and Total Votes"
img_path = "genre_metrics_table.png"
slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), width=Inches(9))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "This table shows the average ratings and total number of votes for each genre in the TMDB dataset."

# Slide 3: Genre Distribution Pie Chart
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Genre Distribution in TMDB Movies Dataset"
img_path = "Grouped_Genre_Distribution_Pie_Chart.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "This pie chart visualizes the distribution of movie genres in the TMDB dataset. The largest genres include Drama, Comedy, and Romance."

# Slide 4: Top Recommendations Visualization
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Top Movie Recommendations"
img_path = "movie_recommendations.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "These are the top movie recommendations based on user preferences and similar movies."

# Slide 5: Preprocessing Workflow
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Data Preprocessing Workflow"
img_path = "preprocessing_workflow.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "This flowchart illustrates the preprocessing steps: merging datasets, extracting features (genres, cast), and normalizing the text."

# Slide 6: AI Model Workflow
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "AI Model Development Workflow"
img_path = "ai_model_workflow_colored.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "This flowchart visualizes the process of model development: inputting a movie, vectorization, similarity calculation, and ranking recommendations."

# Slide 7: Algorithm Comparison
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Algorithm Comparison on RMSE and MAE"
img_path = "algorithm_comparison.png"
slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

# Add supporting text
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
text_frame = textbox.text_frame
text_frame.text = "This chart compares the performance of different algorithms based on RMSE and MAE, showing the accuracy and error metrics for each."

# Save the PowerPoint presentation
output_path = "Movie_Recommendation_System_Presentation_With_Text.pptx"
presentation.save(output_path)

output_path
